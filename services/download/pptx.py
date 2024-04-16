from pptx import Presentation

from config.logger import log


def extract_text(file_path: str) -> str:
    log().info(f"Extracting text from {file_path}")

    prs = Presentation(file_path)
    total_slides = len(prs.slides)

    log().debug(f"Total slides: {total_slides}")

    all_slides = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        texts = []

        log().debug(f"slide: {slide_number} contained {len(slide.shapes)} shapes")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)

        slide_text = '\n'.join(texts)
        slide = f"""
Slide {slide_number}/{total_slides}
- - - - - - - -

{slide_text}
""".strip()
        all_slides.append(slide)

    return "\n\n".join(all_slides)
